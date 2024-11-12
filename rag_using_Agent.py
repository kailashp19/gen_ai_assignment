from qdrant_client import QdrantClient
from qdrant_client.http.models import VectorParams, PointStruct
from pptx import Presentation  # For handling .pptx files
from transformers import BertTokenizer, BertModel
from langchain.tools import Tool
from langchain.agents import initialize_agent, AgentType
import fitz  # PyMuPDF for PDFs
import pytesseract  # For image to text conversion
from PIL import Image  # Python Imaging Library for handling images
import google.generativeai as genai
import os
import re
import docx  # For handling .docx files
import torch
import numpy as np
import warnings
warnings.filterwarnings('ignore')

# Define the Gemini model wrapper to make it compatible
class GeminiModelWrapper(Runnable):
    def __init__(self, model):
        self.model = model

    def __call__(self, prompt):
        response = self.model.generate_content(prompt)
        return response.candidates[0].text if response.candidates else ""

def clean_text(text):
    """Clean and normalize extracted text, removing extra spaces but preserving newlines."""
    text = text.strip()
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    return text

def pdf_to_text(pdf_path, text_path):
    """Convert PDF to text, clean it, and save to a file."""
    doc = fitz.open(pdf_path)
    text = ""
    
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        page_text = page.get_text()
        text += clean_text(page_text) + "\n"

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def image_to_text(image_path, text_path):
    """Convert image to text using pytesseract, clean it, and save to a file."""
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img)
    text = clean_text(text)

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def docx_to_text(docx_path, text_path):
    """Convert DOCX to text, clean it, and save to a file."""
    doc = docx.Document(docx_path)
    text = ""

    for para in doc.paragraphs:
        text += clean_text(para.text) + "\n"

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def pptx_to_text(pptx_path, text_path):
    """Convert PPTX to text, clean it, and save to a file."""
    presentation = Presentation(pptx_path)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += clean_text(shape.text) + "\n"

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def vector_db_creation(file_dir, collection_name):    
    # Load model for generating embeddings
    tokenizer = BertTokenizer.from_pretrained('bert-base-uncased')
    model = BertModel.from_pretrained('bert-base-uncased')

    # define the max length
    max_length = 128

    # Load text from a file
    text_file_path = os.path.join(text_dir, "Error Codes.txt")
    print("Text file path", text_file_path)
    with open("D:/Users/kaila/Personal Projects/Team_6_Gen_AI/gen_ai_capstone/gen_ai_assignment/Error Codes.txt", 'r', encoding='utf-8') as file:
        sentences = [line.strip() for line in file.readlines() if line.strip()]

    # f = open("D:/Users/kaila\Personal Projects/Team_6_Gen_AI/gen_ai_capstone/gen_ai_assignment/Error Codes.txt", 'r')
    # print(f.read())

    # Tokenize and encode sentences
    inputs = tokenizer(sentences, padding='max_length', truncation=True, max_length=max_length, return_tensors="pt")

    # Generate embeddings
    with torch.no_grad():
        outputs = model(**inputs)
        # Use the [CLS] token representation as the sentence embedding
        embeddings = outputs.last_hidden_state[:, 0, :].numpy()
    
    # Save embeddings to a file
    np.save('fixed_length_embeddings.npy', embeddings)

    print("Embeddings generated and saved.")
    print(f"Length of each embedding: {embeddings.shape[0]}")

    vector_size = embeddings.shape[1]  # Should be 768 for 'bert-base-uncased'
    distance_metric = "Cosine"  # Change to "Dot" or "Euclidean" if needed

    # Connect to the Qdrant service
    client = QdrantClient("http://localhost:6333")

    # Create collection in Qdrant
    client.recreate_collection(
        collection_name=collection_name,
        vectors_config=VectorParams(size=vector_size, distance=distance_metric),
    )

    # Prepare and upload embeddings to Qdrant
    points = [
        PointStruct(id=i, vector=embedding.tolist(), payload={"text": sentences[i]})
        for i, embedding in enumerate(embeddings)
    ]

    client.upsert(collection_name=collection_name, points=points)

    # Encode the query into an embedding
    query = "I am getting an error as 1333"
    
    # Tokenize and generate query embedding
    query_inputs = tokenizer(query, padding='max_length', truncation=True, max_length=128, return_tensors="pt")

    with torch.no_grad():
        query_outputs = model(**query_inputs)
        query_embedding = query_outputs.last_hidden_state[:, 0, :].numpy()

    # Perform the search in Qdrant
    search_results = client.search(
        collection_name=collection_name,
        query_vector=query_embedding[0],  # Extract the first vector from the batch
        limit=3  # Retrieve multiple results for context
    )

    # Retrieve the top result text chunks for context
    context_texts = [result.payload["text"] for result in search_results]

    context_prompt = f"""You are provided with a document in .txt format and your task is to provide a solution to a user coming with problem
    for tehcnical issue from a document.
    For Example:
    User query: I am getting an error as 1333
    output: 
    Here is a recommended solution
    Verify the Joey is active on the customer account   
    Perform a front panel reset on the Joey and Hopper
    Ensure the Joey is linked to the Hopper in SETTINGS > WHOLE HOME
    Inspect the signal path starting at the Joey and working back toward the Hub/Node:\n\n"""

    for i, text in enumerate(context_texts, start=1):
        context_prompt += f"Context {i}: {text}\n\n"""

    context_prompt += f"Question: {query}\nAnswer:"

    # Call Gemini API to generate the response
    gemini_response = llm.generate_content(context_prompt)

    # Display the generated response
    print("\nGemini Response:")
    # print(gemini_response["text"])

    if gemini_response.candidates:
        extracted_text = gemini_response.text.strip('```json').strip('```').strip()
        print(f"Extracted text: {extracted_text}")  # Debugging purpose
        return extracted_text
    else:
        extracted_text = ''
        return extracted_text

pdf_tool = Tool(
    name="pdf_to_text",
    func=pdf_to_text,
    description="Converts a PDF file to text."
)

image_tool = Tool(
    name="image_to_text",
    func=image_to_text,
    description="Converts an image file to text."
)

docx_tool = Tool(
    name="docx_to_text",
    func=docx_to_text,
    description="Converts a docx file to text."
)

pptx_tool = Tool(
    name="pptx_to_text",
    func=pptx_to_text,
    description="Converts a pptx file to text."
)

embedding_tool = Tool(
    name="vector_db_creation",
    func=vector_db_creation,
    description="for llm response"
)

tools = [pdf_tool, image_tool, docx_tool, pptx_tool, embedding_tool]

# Configure Google Gemini API
genai.configure(api_key='AIzaSyAt8gpOAHgwzOGOhpJATz88vxMeeM1q2Lg')
llm = genai.GenerativeModel("gemini-1.5-flash")

agent = initialize_agent(
    tools=tools,
    llm=llm,
    agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
    verbose=True
)

def main(file_dir):
    # Process each file in the directory
    for file_name in os.listdir(file_dir):
        print(file_name)
        file_path = os.path.join(file_dir, file_name)
        text_file_name = os.path.splitext(file_name)[0] + '.txt'
        text_path = os.path.join(text_dir, text_file_name)
        if file_name.lower().endswith('.pdf'):
            print(f"Converting {file_path} to {text_path}")
            pdf_tool.func(file_path, text_path)
        elif file_name.lower().endswith(('.jpg', '.jpeg', '.png')):
            print(f"Converting {file_path} to {text_path}")
            image_tool.func(file_path, text_path)
        elif file_name.lower().endswith('.docx'):
            print(f"Converting {file_path} to {text_path}")
            docx_tool.func(file_path, text_path)
        elif file_name.lower().endswith('.pptx'):
            print(f"Converting {file_path} to {text_path}")
            pptx_tool.func(file_path, text_path)
        else:
            print(f"Unsupported file format for {file_path}")

    print("All files have been converted to text files with preserved newlines and no extra spaces.")
    # Initialize the Gemini model
    genai.configure(api_key='YOUR_API_KEY')
    gemini_model = genai.GenerativeModel("gemini-1.5-flash")
    gemini_llm = GeminiModelWrapper(gemini_model)
    response = embedding_tool.func(text_path, collection_name)
    return response

if __name__=="__main__":

    # Directories
    file_dir = "Capstone data sets"  # Directory containing files of various formats
    text_dir = "Converted text files"  # Directory to save text files
    # os.makedirs(text_dir, exist_ok=True)

    # defining the name of the collection
    collection_name = "Technical_Support_Agent"

    result = main(file_dir)
    print("Workflow Result:", result)