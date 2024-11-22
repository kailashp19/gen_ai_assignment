from pptx import Presentation  # For handling .pptx files
from PIL import Image  # Python Imaging Library for handling images
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.http.models import VectorParams, PointStruct
from docx.opc.exceptions import PackageNotFoundError
from nltk import sent_tokenize
from sklearn.metrics.pairwise import cosine_similarity
import fitz  # PyMuPDF for PDFs
import pytesseract  # For image to text conversion
import google.generativeai as genai
import os
import re
import docx  # For handling .docx files
import numpy as np
import warnings
warnings.filterwarnings('ignore')
import nltk

class RAGPRocessing():
    """
    A class which takes file(s) either in .pdf, .docx, .pptx, or image files in jpg or png format which contains the
    information on technical error codes. Once the documents from above format are provided, it then converts the above documents
    into text files. These text files are first broken down into chunks, tokenized it, and then later on embeddings are created out
    of it. After the embedding are created, these embeddings are stored into Vector Database (Qdrant). A user query is also passed 
    to search for the result from the vector database. the searched results are then passed to system prompt and LLM to generate the
    response of the uer query for the technical issues.
    """
    def __init__(self, file_dir: str, text_dir: str, client: object, model: object, collection_name: str, user_query: str, chunk_size: int, llm: str) -> None:
        self.file_dir = file_dir
        self.text_dir = text_dir
        self.client = client
        self.model = model
        self.collection_name = collection_name
        self.user_query = user_query
        self.chunk_size = chunk_size
        self.llm = llm
    def clean_text(self, text: str) -> str:
        """
        Clean and normalize extracted text, removing extra spaces but preserving newlines.

        Parameters:
        text (str): A string to clean and normalize it.

        Return:
        text (str): A string which is cleaned and normalized

        """
        text = text.strip()
        text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
        return text

    def pdf_to_text(self, pdf_path: str, text_path: str) -> None:
        """
        Convert PDF to text, clean it, and save to a file.
        
        Parameters:
        pdf_path (str): A path to pdf file in string format.
        text_path (str): A path to save it as a text file in string format.

        Return:
        None
        """
        doc = fitz.open(pdf_path)
        text = ""
        
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            text += self.clean_text(page_text) + "\n"

        with open(text_path, 'w', encoding='utf-8') as text_file:
            text_file.write(text)

    def image_to_text(self, image_path: str, text_path: str) -> None:
        """
        Convert image to text using pytesseract, clean it, and save to a file.

        Parameters:
        image_path (str): A path to image file in string format.
        text_path (str): A path to save it as a text file in string format.

        Return:
        None

        """
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)
        text = self.clean_text(text)

        with open(text_path, 'w', encoding='utf-8') as text_file:
            text_file.write(text)

    def docx_to_text(self, docx_path: str, text_path: str) -> None:
        """
        Convert DOCX to text, clean it, and save to a file.
        
        Parameters:
        docx_path (str): A path to docx file in string format.
        text_path (str): A path to save it as a text file in string format.

        Return:
        None
        """
        try:
            doc = docx.Document(docx_path)
            text = ""
            for para in doc.paragraphs:
                text += self.clean_text(para.text) + "\n"
            
            with open(text_path, 'w', encoding='utf-8') as text_file:
                text_file.write(text)
        except PackageNotFoundError as e:
            print(f"PackageNotFoundError: {e} - The file {docx_path} is not valid or is corrupted.")
        except Exception as e:
            print(f"Error processing file {docx_path}: {e}")

    def pptx_to_text(self, pptx_path: str, text_path: str) -> None:
        """
        Convert PPTX to text, clean it, and save to a file.
        
        Parameters:
        pptx_path (str): A path to ppt file in string format.
        text_path (str): A path to save it as a text file in string format.

        Return:
        None
        """
        presentation = Presentation(pptx_path)
        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += self.clean_text(shape.text) + "\n"

        with open(text_path, 'w', encoding='utf-8') as text_file:
            text_file.write(text)

    def vector_db_creation(self) -> str:
        """
        Creates the embeddings, store the embeddings into a vector dataabse, search for the user query in the vector database and finally
        passed the search results to an LLM model for more contextualized results.

        Parameters:
        text_dir (str): A directory to text file in string format.
        client (object): A connection string to the Qdrant Client.
        model (object): An embedding model
        collection_name (str): A name of the collection where vector embeddings are stored.
        user_query (str): A user provided query.
        chunk_size (int): A size of the chunk, usually 128 characters
        llm (object): A Gemini 1.5 flash LLM Model

        Returns:
        extracted_text (str): A response from an LLM Model in the form of natural language.
        """
        with open(f"{self.text_dir}/Error Codes.txt", 'r', encoding='utf-8') as file:
            text = file.read()
        
        # Split the document into chunks (sentences in this case)
        sentences = sent_tokenize(text)
        overlap_size = 50
        chunks = []
        for i in range(0, len(sentences), self.chunk_size - overlap_size):
            chunk = ' '.join(sentences[i:i + self.chunk_size])
            chunks.append(chunk)

        # Generate embeddings
        embeddings = self.model.encode(chunks, show_progress_bar=True)

        # Convert to a NumPy array
        embeddings = np.array(embeddings)

        # Save embeddings to a file
        np.save('sentence_embeddings.npy', embeddings)

        # Create a collection to store the embeddings
        self.client.recreate_collection(
            collection_name=self.collection_name,
            vectors_config=VectorParams(size=embeddings.shape[1], distance='Cosine')
        )

        # Add points to the collection
        points = [
            PointStruct(id=i, vector=embedding, payload={"document": doc})
            for i, (embedding, doc) in enumerate(zip(embeddings, chunks))
        ]

        self.client.upsert(collection_name=self.collection_name, points=points)

        # Encode the query into an embedding
        query_embedding = self.model.encode([self.user_query])[0]  # Flatten the query embedding

        # Perform the search
        search_results = self.client.search(
            collection_name=self.collection_name,
            query_vector=query_embedding,
            limit=5  # Number of nearest neighbors to return
        )

        # Display the search results
        for idx, result in enumerate(search_results, start=1):
            text_chunk = result.payload.get("document", "No text found")
            score = result.score
            # print(f"accuracy score: {score}")

        # Retrieve the top result text chunks for context
        context_texts = search_results

        context_prompt = f"""
        You are a technical assistant and your job is to accurately answer the user query.
        You are provided with a document in .txt format and your task is to answer user query specific to search result only.

        Remember:
        1. *Accuracy is crucial* for the user query.
        2. Provide reponse specific to user query. i.e if asked for the solution, then respond with solution. If asked for cause,
        provide them with exact cause so that user can understand it clearly. If asked for meaning, then provide the full detail of the error.
        3. Provide the response in user-friendly way.

        EXAMPLES:
        USER QUERY: I am getting an error as 1333
        RESPONSE: Error 1333 says that the device is looking for hopper.

        USER QUERY: What does 1333 means?
        RESPONSE: Sorry please provide more context what does 1333 means?

        USER QUERY: I am getting an error code 002, what should I do?
        RESPONSE: Error code 002 could be caused by any of the following:
        1. Signal loss due to inclement weather. 
        2. Misaligned or damaged dish antenna. 
        3. Failure of component
        4. Improperly installed component in signal path.
        Try any of the below resolutions:
        1. Using a signal meter, check signal at the receiver location and work back toward the dish antenna.
        2. Check signal at each connection point.
        3. Examine coax cable, connectors and components for proper installation and approved parts. 
        4. Correct signal issue at dish antenna
        
        Context
        {context_texts}
        User Query
        {self.user_query}\n\n"""

        # Call Gemini API to generate the response
        gemini_response = self.llm.generate_content(context_prompt)
        if gemini_response.candidates:
            extracted_text = gemini_response.text.strip('```json').strip('```').strip()
            return extracted_text
        else:
            extracted_text = ''
            return extracted_text
        
    def calculate_rag_metrics(self, response):
        """
        A method to calculate various RAG metrics from the response generated by LLM.

        Parameter:
        response (str): generated response from LLM.

        Reurns:
        rag_matrix (dict): rag evaluation matrix.
        """

        with open(f"{self.text_dir}/Error Codes.txt", 'r', encoding='utf-8') as file:
            text = file.read()
        
        # Split the document into chunks (sentences in this case)
        sentences = sent_tokenize(text)
        response_sentences = sent_tokenize(response)

        chunks = [' '.join(sentences[i:i + self.chunk_size]) for i in range(0, len(sentences), self.chunk_size)]
        response_chunks = [' '.join(response_sentences[i:i + self.chunk_size]) for i in range(0, len(response_sentences), self.chunk_size)]
        query_chunks = [' '.join(self.user_query[i:i + self.chunk_size]) for i in range(0, len(self.user_query), self.chunk_size)]

        # Generate embeddings
        embeddings = self.model.encode(chunks, show_progress_bar=True)

        # Convert to a NumPy array
        embeddings = np.array(embeddings)

        # Encode the query into an embedding
        query_embedding = self.model.encode(query_chunks, show_progress_bar=True)  # Flatten the query embedding
        response_embeddings = self.model.encode(response_chunks, show_progress_bar=True)

        # query_embedding = np.array(query_embedding)
        response_embeddings = np.array(response_embeddings)

        # print(response_embeddings.shape)
        # print(embeddings.shape)
        # print(query_embedding.shape)

        if embeddings.shape[0] == 1:
            retrieved_embeddings = embeddings.reshape(1, -1)

        if response_embeddings.shape[0] == 1:
            response_embeddings = response_embeddings.reshape(1, -1)

        # calculating cosine similarity score
        score = cosine_similarity(response_embeddings, retrieved_embeddings)
        # print("Cosine Similarity: ", score)

        response_tokens = set(response.split())
        context_tokens = set(' '.join(chunks).split())
        # Precision: Proportion of response tokens found in the context
        precision = len(response_tokens & context_tokens) / len(response_tokens)
        # Recall: Proportion of context tokens found in the response
        recall = len(response_tokens & context_tokens) / len(context_tokens)
        # print("Precision: ", precision)
        # print("Recall: ", recall)

        relevancy_score = cosine_similarity(query_embedding, response_embeddings)[0]
        # print("Answer relevancy score: ", relevancy_score)
        rag_matrix = {
            "cosine similarity": score,
            "precision": precision,
            "recall": recall,
            "relevancy score": relevancy_score
        }
        return rag_matrix

def main(file_dir, text_dir, client, model, collection_name, user_query, chunk_size, llm):
    pdoc = RAGPRocessing(file_dir, text_dir, client, model, collection_name, user_query, chunk_size, llm)
    # Process each file in the directory
    for file_name in os.listdir(file_dir):
        file_path = os.path.join(file_dir, file_name)
        text_file_name = os.path.splitext(file_name)[0] + '.txt'
        text_path = os.path.join(text_dir, text_file_name)
        if file_name.lower().endswith('.pdf'):
            pdoc.pdf_to_text(file_path, text_path)
        elif file_name.lower().endswith(('.jpg', '.jpeg', '.png')):
            pdoc.image_to_text(file_path, text_path)
        elif file_name.lower().endswith('.docx'):
            pdoc.docx_to_text(file_path, text_path)
        elif file_name.lower().endswith('.pptx'):
            pdoc.pptx_to_text(file_path, text_path)
        else:
            print(f"Unsupported file format for {file_path}")
    result = pdoc.vector_db_creation()
    rag_matrix = pdoc.calculate_rag_metrics(result)
    return result

if __name__=="__main__":
    # Directories
    file_dir = "Capstone data sets"  # Directory containing files of various formats
    text_dir = "Converted text files"  # Directory to save text files
    os.makedirs(text_dir, exist_ok=True)
    # Connect to the Qdrant service
    client = QdrantClient("http://localhost:6333")
    model = SentenceTransformer('BAAI/bge-large-en-v1.5')
    # defining the name of the collection
    collection_name = "Technical_Support_Agent"
    # Configure Google Gemini API
    genai.configure(api_key='sk-')
    llm = genai.GenerativeModel('gemini-1.5-flash')
    user_query = "What does error code 002 means?"
    chunk_size = 512
    res = main(file_dir, text_dir, client, model, collection_name, user_query, chunk_size, llm)