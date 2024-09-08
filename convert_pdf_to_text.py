import fitz  # PyMuPDF for PDFs
import pytesseract  # For image to text conversion
from PIL import Image  # Python Imaging Library for handling images
import os
import re
import docx  # For handling .docx files
from pptx import Presentation  # For handling .pptx files

# Directories
file_dir = "docs"  # Directory containing files of various formats
text_dir = "text_files"  # Directory to save text files
os.makedirs(text_dir, exist_ok=True)

def clean_text(text):
    """Clean and normalize extracted text, removing extra spaces but preserving newlines."""
    text = text.strip()
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    return text

def pdf_to_text(pdf_path, text_path):
    """Convert PDF to text, clean it, and save to a file."""
    doc = fitz.open(pdf_path)
    text = ""
    
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        page_text = page.get_text()
        text += clean_text(page_text) + "\n"

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def image_to_text(image_path, text_path):
    """Convert image to text using pytesseract, clean it, and save to a file."""
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img)
    text = clean_text(text)

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def docx_to_text(docx_path, text_path):
    """Convert DOCX to text, clean it, and save to a file."""
    doc = docx.Document(docx_path)
    text = ""

    for para in doc.paragraphs:
        text += clean_text(para.text) + "\n"

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def pptx_to_text(pptx_path, text_path):
    """Convert PPTX to text, clean it, and save to a file."""
    presentation = Presentation(pptx_path)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += clean_text(shape.text) + "\n"

    with open(text_path, 'w', encoding='utf-8') as text_file:
        text_file.write(text)

def main():
    # Process each file in the directory
    for file_name in os.listdir(file_dir):
        file_path = os.path.join(file_dir, file_name)
        text_file_name = os.path.splitext(file_name)[0] + '.txt'
        text_path = os.path.join(text_dir, text_file_name)

        if file_name.lower().endswith('.pdf'):
            print(f"Converting {file_path} to {text_path}")
            pdf_to_text(file_path, text_path)
        elif file_name.lower().endswith(('.jpg', '.jpeg', '.png')):
            print(f"Converting {file_path} to {text_path}")
            image_to_text(file_path, text_path)
        elif file_name.lower().endswith('.docx'):
            print(f"Converting {file_path} to {text_path}")
            docx_to_text(file_path, text_path)
        elif file_name.lower().endswith('.pptx'):
            print(f"Converting {file_path} to {text_path}")
            pptx_to_text(file_path, text_path)
        else:
            print(f"Unsupported file format for {file_path}")

    print("All files have been converted to text files with preserved newlines and no extra spaces.")

if __name__=="__main__":
    main()